#!/usr/bin/env python3
"""Self-Attentionの仕組みを解説するパワーポイント (v3 デザイン刷新)"""

import os, matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 洗練カラーパレット (色数抑えめ) ──
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_GRAY = RGBColor(0xE9, 0xEC, 0xEF)
MID_GRAY   = RGBColor(0x6C, 0x75, 0x7D)
DARK_GRAY  = RGBColor(0x34, 0x3A, 0x40)
NEAR_BLACK = RGBColor(0x21, 0x25, 0x29)
PRIMARY    = RGBColor(0x0D, 0x6E, 0xFD)  # 青
PRIMARY_DK = RGBColor(0x08, 0x4E, 0xB5)
ACCENT     = RGBColor(0x19, 0x87, 0x54)  # 緑
WARN       = RGBColor(0xFD, 0x7E, 0x14)  # オレンジ
KEY_BG     = RGBColor(0xE7, 0xF1, 0xFF)  # 薄い青
CODE_BG    = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG    = RGBColor(0xCD, 0xD6, 0xF4)

MONO = "Consolas"
SANS = "Meiryo"

EQ_DIR = "/tmp/sa_equations_v3"
os.makedirs(EQ_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

_eq_counter = 0
_slide_num  = 0

# ── 数式レンダリング ──
def render_eq(latex, fontsize=24, figw=8, figh=0.8, color="black"):
    global _eq_counter; _eq_counter += 1
    path = os.path.join(EQ_DIR, f"eq_{_eq_counter:03d}.png")
    fig, ax = plt.subplots(figsize=(figw, figh))
    ax.text(0.5, 0.5, f"${latex}$", fontsize=fontsize, ha='center', va='center',
            transform=ax.transAxes, color=color, math_fontfamily='cm')
    ax.axis('off')
    fig.savefig(path, dpi=200, bbox_inches='tight', transparent=True, pad_inches=0.05)
    plt.close(fig); return path

def add_eq(slide, latex, left, top, width, height, fontsize=24):
    path = render_eq(latex, fontsize, figw=width*1.2, figh=height*1.2)
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))

# ── 共通ヘルパー ──
def _tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def _run(para, text, sz=18, bold=False, italic=False, color=DARK_GRAY, name=SANS):
    r = para.add_run(); r.text = text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name=name; return r

def _para(tf, text="", sz=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
          sb=Pt(4), sa=Pt(2), name=SANS):
    p = tf.add_paragraph(); p.alignment=align; p.space_before=sb; p.space_after=sa
    if text: _run(p, text, sz, bold, color=color, name=name)
    return p

def _box(slide, l, t, w, h, fill=KEY_BG, border=None, radius=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    return sh

def _code(slide, l, t, w, h, text, fsz=14):
    sh = _box(slide, l, t, w, h, fill=CODE_BG)
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left=Inches(0.35); tf.margin_right=Inches(0.35)
    tf.margin_top=Inches(0.2); tf.margin_bottom=Inches(0.2)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(1); p.space_after=Pt(1)
        _run(p, line, sz=fsz, color=CODE_FG, name=MONO)
    return sh

def _bullets(tf, items, sz=18, color=DARK_GRAY):
    for item in items:
        p = tf.add_paragraph(); p.space_before=Pt(8); p.space_after=Pt(2)
        if isinstance(item, tuple):
            _run(p, item[0], sz=sz, bold=True, color=DARK_GRAY)
            _run(p, item[1], sz=sz, color=color)
        else:
            _run(p, item, sz=sz, color=color)

# ── スライドテンプレート ──
def _title_slide(title, subtitle=""):
    global _slide_num; _slide_num = 0
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = NEAR_BLACK
    # アクセントライン
    line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.1), Inches(2.5), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = PRIMARY; line.line.fill.background()
    tb = _tb(sl, 1.5, 1.6, 10.3, 1.5)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _run(p, title, sz=44, bold=True, color=WHITE)
    if subtitle:
        _para(tf, subtitle, sz=20, color=MID_GRAY, sb=Pt(24))
    return sl

def _content(title):
    global _slide_num; _slide_num += 1
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_WHITE
    # タイトル帯 (左端にアクセント色の縦線)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Pt(5), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
    tb = _tb(sl, 0.5, 0.15, 12, 0.7)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _run(p, title, sz=26, bold=True, color=NEAR_BLACK)
    # スライド番号
    num = _tb(sl, 12.3, 7.0, 0.8, 0.4)
    ntf = num.text_frame
    p = ntf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, str(_slide_num), sz=12, color=MID_GRAY)
    # 薄い区切り線
    sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12.3), Pt(1))
    sep.fill.solid(); sep.fill.fore_color.rgb = LIGHT_GRAY; sep.line.fill.background()
    return sl

def _keypoint(sl, l, t, w, h, text, sz=18):
    """重要ポイント用テキスト（背景なし）"""
    tb = _tb(sl, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _run(p, text, sz=sz, bold=True, color=PRIMARY_DK)
    return tb

def _circle_num(sl, cx, cy, num, color=PRIMARY):
    """番号付き丸を描画"""
    sz = 0.5
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx - sz/2), Inches(cy - sz/2), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, str(num), sz=18, bold=True, color=WHITE)

def _arrow(sl, x, y):
    """矢印テキスト"""
    tb = _tb(sl, x, y-0.15, 0.5, 0.3)
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "→", sz=22, color=MID_GRAY)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド作成
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# ── 1. タイトル ──
_title_slide("Self-Attentionの仕組み", "直観と数式と手計算で理解する")

# ── 2. 核心 ──
s = _content("Self-Attentionがやっていることの核心")

tb = _tb(s, 0.8, 1.2, 11.5, 1.2)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "各トークンが「列全体のどこを参照すべきか」を自分で決め、参照結果を取り込んで", sz=19, color=DARK_GRAY)
_run(p, "文脈化された表現", sz=19, bold=True, color=PRIMARY)
_run(p, "に更新する仕組み。", sz=19, color=DARK_GRAY)

add_eq(s,
    r"\mathrm{Attention}(Q,K,V) = \mathrm{softmax}\!\left(\frac{QK^\top}{\sqrt{d_k}}\right)V",
    2.8, 2.5, 7.7, 1.0, fontsize=28)

# 3ステップ: 縦並びで各ステップを詳しく説明
steps = [
    ("類似度スコアを計算",
     "各トークンの Query と、全トークンの Key の内積を取る。"
     "内積が大きいほど「方向が揃っている＝相性が良い」ことを意味する。"),
    ("Softmax で重みに変換",
     "内積スコアを sqrt(d_k) で割ってから softmax を適用。"
     "全トークンへの重みが 0〜1 の確率分布になる（和 = 1）。"),
    ("Value の加重和で出力を作る",
     "重みに従って全トークンの Value をブレンド（加重和）。"
     "結果が「文脈を取り込んだ新しい表現」になる。"),
]
y_start = 3.8
for i, (label, desc) in enumerate(steps):
    y = y_start + i * 1.2
    _circle_num(s, 1.3, y + 0.25, i+1)
    tb = _tb(s, 1.7, y, 10.5, 0.4)
    tf = tb.text_frame; p = tf.paragraphs[0]
    _run(p, label, sz=18, bold=True, color=NEAR_BLACK)
    tb2 = _tb(s, 1.7, y + 0.4, 10.5, 0.6)
    tf2 = tb2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]
    _run(p2, desc, sz=15, color=MID_GRAY)

# ── 2b. Self-Attention が生まれた背景 ──
s = _content("Self-Attention が生まれるまで — RNN/LSTMの限界")

# RNN時代の問題を3列で
problems = [
    ("RNN (2014以前)",
     "単語を1つずつ順番に処理",
     "逐次処理 → GPU並列化できず\n学習が遅い",
     WARN),
    ("LSTM / GRU",
     "ゲート機構で長期記憶を改善",
     "逐次処理は変わらず\n100語超で依存関係が薄れる",
     WARN),
    ("Attention (2014)",
     "Bahdanau: RNNの出力に\nAttentionを追加",
     "遠い単語も直接参照できた\nただしRNNの上に載せる形",
     ACCENT),
]
for i, (name, what, limit, color) in enumerate(problems):
    x = 0.5 + i * 4.2
    sh = _box(s, x, 1.2, 3.8, 3.0, fill=WHITE, border=color)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.15)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, name, sz=17, bold=True, color=color)
    _para(tf, what, sz=14, color=DARK_GRAY, align=PP_ALIGN.CENTER, sb=Pt(12))
    _para(tf, "", sz=6)
    _para(tf, limit, sz=13, color=MID_GRAY, align=PP_ALIGN.CENTER, sb=Pt(6))

# 転換点
tb = _tb(s, 0.8, 4.5, 11.5, 0.5)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "2017 — \"Attention Is All You Need\" (Vaswani et al.)", sz=20, bold=True, color=PRIMARY)

tb2 = _tb(s, 0.8, 5.1, 11.5, 1.8)
tf2 = tb2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_run(p, "気づき: ", sz=18, bold=True, color=NEAR_BLACK)
_run(p, "Bahdanau Attentionの実験で、長距離の依存関係を捉えていたのは", sz=18, color=DARK_GRAY)
_run(p, "RNNではなくAttention部分", sz=18, bold=True, color=NEAR_BLACK)
_run(p, "だった", sz=18, color=DARK_GRAY)
p2 = _para(tf2, "", sz=18, sb=Pt(8))
_run(p2, "→ RNNはもはやボトルネック。", sz=18, color=DARK_GRAY)
_run(p2, "Attentionだけで系列を処理すればいい", sz=18, bold=True, color=NEAR_BLACK)
p3 = _para(tf2, "→ 全トークンが全トークンを直接参照（逐次処理なし、行列積で一括計算）", sz=18, color=DARK_GRAY, sb=Pt(8))

# ── 2c. Self-Attentionは何をするのか ──
s = _content("Self-Attention は何をするのか")

tb = _tb(s, 0.8, 1.2, 11.5, 0.5)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "問題: 単語の埋め込みは「文脈を知らない」", sz=21, bold=True, color=PRIMARY)

_code(s, 0.8, 1.8, 11.5, 1.9,
    '  "neko" -> x_neko = [0.2, 0.8, ...]    always the same vector!\n'
    "\n"
    '  "neko ga sakana wo tabeta"    <- neko = subject (doer)\n'
    '  "inu  ga neko   wo oikaketa"  <- neko = object  (chased)', fsz=15)

tb2 = _tb(s, 0.8, 3.9, 11.5, 0.8)
tf2 = tb2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_run(p, "「猫」の役割は周囲の単語で決まる。", sz=19, color=DARK_GRAY)
_run(p, "でも埋め込みベクトルは常に同じ。", sz=19, bold=True, color=NEAR_BLACK)
_para(tf2, "→ 周囲の単語の情報を取り込んで、文脈に合った表現に更新したい", sz=19, color=DARK_GRAY, sb=Pt(8))

tb3 = _tb(s, 0.8, 5.0, 11.5, 0.4)
tf3 = tb3.text_frame; p = tf3.paragraphs[0]
_run(p, "でも「全員の情報を均等に混ぜる」のではダメ:", sz=18, bold=True, color=NEAR_BLACK)

tb4 = _tb(s, 1.3, 5.5, 11.0, 0.8)
tf4 = tb4.text_frame; tf4.word_wrap = True
p = tf4.paragraphs[0]
_run(p, "「猫が魚を食べた」で「食べた」を更新するとき、", sz=17, color=DARK_GRAY)
_run(p, "「猫」(主語)と「を」(助詞)では重要度が全然違う。", sz=17, bold=True, color=DARK_GRAY)
_para(tf4, "→ トークンごとに「誰がどれくらい重要か」をスコア化する仕組みが要る", sz=17, color=DARK_GRAY, sb=Pt(6))

_keypoint(s, 0.8, 6.6, 11.5, 0.6,
    "Self-Attention = 各トークンが「誰の情報をどれだけ取り込むか」のスコアを計算し、重み付きで混ぜる仕組み")

# ── 3. Q, K, V ──
s = _content("Query / Key / Value の直感的な意味")

# まず入力xとは何かを具体的に
tb = _tb(s, 0.8, 1.1, 11.5, 0.4)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "入力 x とは: 各トークン(単語)を表す数値ベクトル(埋め込み)", sz=17, bold=True, color=NEAR_BLACK)

_code(s, 0.8, 1.6, 11.5, 1.5,
    "input:  [neko]  [ga]  [sakana]  [wo]  [tabeta]\n"
    "\n"
    "each token -> embedding vector (d-dim, learned):\n"
    '  x_neko = [0.2, 0.8, ...]   x_tabeta = [0.5, 0.4, ...]   ...', fsz=15)

# 各xから Q, K, V を作る
tb2 = _tb(s, 0.8, 3.3, 11.5, 0.4)
tf2 = tb2.text_frame; p = tf2.paragraphs[0]
_run(p, "各トークンの x に重み行列を掛けて、3つの役割に変換:", sz=17, bold=True, color=NEAR_BLACK)

_code(s, 0.8, 3.8, 11.5, 2.2,
    "x_tabeta * W_Q -> q_tabeta   Q: \"who did the eating? what was eaten?\"\n"
    "x_tabeta * W_K -> k_tabeta   K: \"I am: a verb, action\"\n"
    "x_tabeta * W_V -> v_tabeta   V: (meaning of tabeta to share)\n"
    "\n"
    "x_neko   * W_Q -> q_neko     Q: \"what action applies to me?\"\n"
    "x_neko   * W_K -> k_neko     K: \"I am: a noun, subject\"\n"
    "x_neko   * W_V -> v_neko     V: (meaning of neko to share)", fsz=14)

add_eq(s,
    r"q_i = x_i \, W_Q \qquad k_i = x_i \, W_K \qquad v_i = x_i \, W_V",
    2.5, 6.2, 8.3, 0.5, fontsize=22)

_keypoint(s, 0.8, 6.8, 11.5, 0.5,
    "q_tabeta と k_neko の内積が大きい → 「食べた」は「猫」の情報(v_neko)を多く取り込む")

# ── 4. なぜ内積なのか (比較) ──
s = _content("類似度スコアはなぜ「内積」なのか")

tb = _tb(s, 0.8, 1.2, 11.5, 0.6)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "やりたいこと: ", sz=19, bold=True, color=NEAR_BLACK)
_run(p, "2つのベクトル(q, k)の「相性」を", sz=19, color=DARK_GRAY)
_run(p, "1つの数字", sz=19, bold=True, color=PRIMARY)
_run(p, "にしたい", sz=19, color=DARK_GRAY)

methods = [
    ("ユークリッド距離",  "||q - k||",  "方向が同じでも長さが違うと遠くなる",           False),
    ("コサイン類似度",    "cos(θ)",     "方向だけ見る。良いが正規化が一手間",            False),
    ("内積  q · k",       "方向+長さ→スコア", "行列積 QK^T で全ペア一括計算!",         True),
    ("学習したMLP",       "MLP(q, k)",  "表現力は高いが計算が重すぎる",                 False),
]
for i, (name, what, note, selected) in enumerate(methods):
    x = 0.5 + i * 3.2
    y = 2.1
    bg = KEY_BG if selected else WHITE
    bc = PRIMARY if selected else LIGHT_GRAY
    sh = _box(s, x, y, 2.9, 2.6, fill=bg, border=bc)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.15)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, name, sz=15, bold=True, color=PRIMARY if selected else DARK_GRAY)
    _para(tf, what, sz=13, color=DARK_GRAY, align=PP_ALIGN.CENTER, sb=Pt(10))
    _para(tf, "", sz=4)
    color = ACCENT if selected else MID_GRAY
    _para(tf, note, sz=12, bold=selected, color=color, align=PP_ALIGN.CENTER, sb=Pt(6))

_keypoint(s, 0.8, 5.1, 11.5, 0.7,
    "Transformerが内積を選んだ最大の理由: QK^T の1回の行列積で全トークンペアのスコアが同時に出る")

# ── 5. 内積の幾何 ──
s = _content("内積 = 「ベクトルの方向が揃っているか」を測る")

add_eq(s,
    r"\vec{a} \cdot \vec{b} \;=\; |\vec{a}|\;|\vec{b}|\;\cos\theta",
    3.5, 1.2, 6.3, 0.85, fontsize=30)

tb = _tb(s, 0.8, 2.2, 11.5, 0.5)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "同じ方向 → 大きい正　　直角 → 0　　逆方向 → 負", sz=18, color=DARK_GRAY)

tb2 = _tb(s, 0.8, 2.9, 4, 0.4)
tf2 = tb2.text_frame; p = tf2.paragraphs[0]
_run(p, "2次元の具体例", sz=18, bold=True, color=NEAR_BLACK)

_code(s, 0.8, 3.4, 11.5, 2.3,
    'q  = [ 3, 0]                                  (right-pointing vector)\n'
    '\n'
    'k1 = [ 2, 0]   q.k1 = 3* 2 + 0*0 =  6        same direction  -> large +\n'
    'k2 = [ 0, 2]   q.k2 = 3* 0 + 0*2 =  0        orthogonal      -> zero\n'
    'k3 = [-1, 0]   q.k3 = 3*-1 + 0*0 = -3        opposite        -> negative', fsz=16)

_keypoint(s, 0.8, 6.0, 11.5, 0.7,
    "内積そのものは「方向の一致度を測る」だけの単純な演算")

# ── 6. 学習が意味を与える ──
s = _content("学習(W_Q, W_K)が内積に意味のある相性を与える")

tb = _tb(s, 0.8, 1.2, 11.5, 1.0)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "元の埋め込みをそのまま内積しても意味のある「相性」にはならない。", sz=19, color=DARK_GRAY)
p2 = _para(tf, "", sz=19)
_run(p2, "だから W_Q と W_K で「内積を取る前に」ベクトルを変換する。", sz=19, bold=True, color=NEAR_BLACK)

_code(s, 2.0, 2.5, 9.3, 2.3,
    "  x_i                  x_j\n"
    "   |                    |\n"
    "   | * W_Q              | * W_K\n"
    "   v                    v\n"
    "  q_i                  k_j\n"
    "            \\       /\n"
    "         q_i . k_j  =  score", fsz=16)

tb2 = _tb(s, 0.8, 5.1, 11.5, 0.4)
tf2 = tb2.text_frame; p = tf2.paragraphs[0]
_run(p, "まとめ: 類似度スコアと内積の間にあるもの", sz=18, bold=True, color=NEAR_BLACK)

items = [
    ("① 目的: ", "2つのベクトルの相性を1つの数字にしたい"),
    ("② 性質: ", "内積は「方向の一致度」を測れる"),
    ("③ 効率: ", "行列積で全ペア一括計算 (GPU最適)"),
    ("④ 学習: ", "W_Q, W_K で事前に変換する → 生の埋め込みの内積ではない"),
]
tb3 = _tb(s, 0.8, 5.6, 11.5, 1.5)
tf3 = tb3.text_frame; tf3.word_wrap = True
for label, desc in items:
    p = tf3.add_paragraph(); p.space_before = Pt(3); p.space_after = Pt(1)
    _run(p, label, sz=16, bold=True, color=PRIMARY)
    _run(p, desc, sz=16, color=DARK_GRAY)

# ── 7. sqrt(d_k) ──
s = _content("なぜ sqrt(d_k) で割るのか")

tb = _tb(s, 0.8, 1.2, 11.5, 0.5)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "内積は各成分の積の合計:", sz=19, color=DARK_GRAY)

add_eq(s, r"q \cdot k = \sum_{i=1}^{d_k} q_i \, k_i", 1.5, 1.9, 5.0, 0.8, fontsize=26)
add_eq(s, r"\mathrm{Var}(q \cdot k) \propto d_k", 7.5, 1.9, 4.5, 0.8, fontsize=26)

tb2 = _tb(s, 0.8, 3.0, 11.5, 1.8)
tf2 = tb2.text_frame; tf2.word_wrap = True
_bullets(tf2, [
    ("d_k が増える → ", "足し算の項が増え、内積の値が大きくなる"),
    ("内積が巨大 → ", "softmaxが飽和 (ほぼ0か1) → 勾配が消える"),
], sz=18)

add_eq(s, r"\mathrm{score} = \frac{q_i \cdot k_j}{\sqrt{d_k}}", 3.5, 4.8, 6.0, 1.0, fontsize=28)

_keypoint(s, 0.8, 6.1, 11.5, 0.6,
    "分散を1に戻してsoftmaxの飽和を防ぎ、学習を安定させる")

# ── 8. softmax + 加重和 ──
s = _content("Softmax → 加重和 で「情報を取り込む」")

# 左カラム: softmax
tb = _tb(s, 0.8, 1.2, 5.5, 0.4)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "Softmaxで重みにする理由", sz=20, bold=True, color=NEAR_BLACK)

tb2 = _tb(s, 0.8, 1.7, 5.5, 2.3)
tf2 = tb2.text_frame; tf2.word_wrap = True
_bullets(tf2, [
    ("重みが正 & 和が1 → ", "Valueの凸結合、スケール安定"),
    ("勝者を強調 → ", "指数で差が広がり鋭く選べる"),
    ("微分可能 → ", "誤差逆伝播で学習できる"),
], sz=17)

# 右カラム: 加重和
tb3 = _tb(s, 7.0, 1.2, 5.5, 0.4)
tf3 = tb3.text_frame; p = tf3.paragraphs[0]
_run(p, "加重和 = 情報の取り込み", sz=20, bold=True, color=NEAR_BLACK)

add_eq(s, r"y_i = \sum_j \alpha_{ij} \; v_j", 7.5, 1.8, 4.5, 0.7, fontsize=26)

tb4 = _tb(s, 7.0, 2.6, 5.5, 1.5)
tf4 = tb4.text_frame; tf4.word_wrap = True
_bullets(tf4, [
    ("a_ij が大きいほど ", "トークン j の情報が強く混ざる"),
    ("y_i = ", "列全体を見た上での要約"),
], sz=17)

_keypoint(s, 0.8, 4.3, 11.5, 0.7,
    "本質: 「どの相手から情報をコピー/ブレンドするか」を学習する")

# ── 9. 全体の流れ ──
s = _content("Self-Attention 全体の流れ")

_code(s, 1.5, 1.2, 10.3, 5.5,
    "token embedding  x_i\n"
    "    |\n"
    "    +-- x * W_Q --> q_i        (Query)\n"
    "    +-- x * W_K --> k_i        (Key)\n"
    "    +-- x * W_V --> v_i        (Value)\n"
    "\n"
    "    q_i . k_j  =  score        ... (1) similarity\n"
    "        |\n"
    "        / sqrt(d_k)            ... scaling\n"
    "        |\n"
    "        softmax                ... (2) weights (0~1, sum=1)\n"
    "        |\n"
    "        sum( a_ij * v_j )      ... (3) weighted sum of Values\n"
    "        |\n"
    "        v\n"
    "        y_i  =  context-aware representation", fsz=17)

# ── 9b. 全トークンで同時に起きる ──
s = _content("これが全トークンで同時に起きる")

tb = _tb(s, 0.8, 1.2, 11.5, 0.8)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "前スライドの流れは「1つのトークンを更新する」手順だった。", sz=18, color=DARK_GRAY)
_para(tf, "Self-Attention では全トークンが同時に Q, K, V を出し、全員が全員を参照する:", sz=18, bold=True, color=NEAR_BLACK, sb=Pt(6))

_code(s, 0.8, 2.4, 11.5, 2.0,
    "each token produces all three:\n"
    "  A -> q_A, k_A, v_A\n"
    "  B -> q_B, k_B, v_B\n"
    "  C -> q_C, k_C, v_C", fsz=16)

_code(s, 0.8, 4.6, 11.5, 1.9,
    "new A = q_A vs [k_A, k_B, k_C] -> weighted sum of [v_A, v_B, v_C]\n"
    "new B = q_B vs [k_A, k_B, k_C] -> weighted sum of [v_A, v_B, v_C]\n"
    "new C = q_C vs [k_A, k_B, k_C] -> weighted sum of [v_A, v_B, v_C]\n"
    "                ^^^^^^^^^^^^^^^^^^^   all Keys/Values are shared", fsz=16)

_keypoint(s, 0.8, 6.7, 11.5, 0.5,
    "行列積 QK^T ひとつで全ペアのスコアが同時に計算される")

# ── 10. 手計算: 設定 ──
s = _content("手計算: 入力と重み行列の設定")

tb = _tb(s, 0.8, 1.2, 11.5, 0.4)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "3トークン, 2次元で全要素を明示して最後まで計算", sz=17, color=MID_GRAY)

_code(s, 0.5, 1.8, 12.3, 2.0,
    "Input X (3x2):         Weight matrices (2x2 each, learned):\n"
    "\n"
    "  |1  0|  <- tok1        W_Q = |1 0|    W_K = |1 1|    W_V = |2 0|\n"
    "  |0  1|  <- tok2              |1 1|          |0 1|          |0 1|\n"
    "  |1  1|  <- tok3", fsz=16)

tb2 = _tb(s, 0.8, 4.0, 3, 0.4)
tf2 = tb2.text_frame; p = tf2.paragraphs[0]
_run(p, "Step 1: Q, K, V を計算", sz=19, bold=True, color=NEAR_BLACK)

add_eq(s, r"Q = XW_Q", 0.8, 4.5, 2.2, 0.45, fontsize=20)
add_eq(s, r"K = XW_K", 4.5, 4.5, 2.2, 0.45, fontsize=20)
add_eq(s, r"V = XW_V", 8.2, 4.5, 2.2, 0.45, fontsize=20)

_code(s, 0.8, 5.1, 3.0, 1.2, "  |1  0|\n  |1  1|\n  |2  1|", fsz=16)
_code(s, 4.5, 5.1, 3.0, 1.2, "  |1  1|\n  |0  1|\n  |1  2|", fsz=16)
_code(s, 8.2, 5.1, 3.0, 1.2, "  |2  0|\n  |0  1|\n  |2  1|", fsz=16)

# ── 11. スコア+スケーリング ──
s = _content("手計算: スコア行列とスケーリング")

tb = _tb(s, 0.8, 1.2, 8, 0.4)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "Step 2: 内積で相性を測る", sz=19, bold=True, color=NEAR_BLACK)

add_eq(s, r"S = QK^\top", 0.8, 1.7, 2.5, 0.5, fontsize=22)

_code(s, 3.8, 1.7, 8.5, 1.5,
    "S = |1 0| x |1 0 1| = |1  0  1|  tok1: high w/ 1,3\n"
    "    |1 1|   |1 1 2|   |2  1  3|  tok2: highest at 3\n"
    "    |2 1|             |3  1  4|  tok3: highest at self", fsz=15)

tb2 = _tb(s, 0.8, 3.5, 11.5, 0.4)
tf2 = tb2.text_frame; p = tf2.paragraphs[0]
_run(p, "S_ij = トークン i の Query と トークン j の Key の内積", sz=16, color=MID_GRAY)

tb3 = _tb(s, 0.8, 4.1, 8, 0.4)
tf3 = tb3.text_frame; p = tf3.paragraphs[0]
_run(p, "Step 3: スケーリング", sz=19, bold=True, color=NEAR_BLACK)

add_eq(s, r"\tilde{S} = \frac{S}{\sqrt{2}}", 0.8, 4.6, 3.5, 0.9, fontsize=22)

_code(s, 5.0, 4.6, 7.3, 1.2,
    " |0.707   0       0.707|\n"
    " |1.414   0.707   2.121|\n"
    " |2.121   0.707   2.828|", fsz=16)

# ── 12. softmax ──
s = _content("手計算: Softmax → 重み行列 A")

tb = _tb(s, 0.8, 1.2, 10, 0.4)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "Step 4: 行ごとにsoftmax (各Queryごとに重み分布を作る)", sz=19, bold=True, color=NEAR_BLACK)

_code(s, 0.8, 1.8, 11.5, 3.3,
    "row 1: [0.707, 0, 0.707]\n"
    "  -> exp: [2.028, 1.000, 2.028]  -> normalize -> [0.401, 0.198, 0.401]\n"
    "\n"
    "row 2: [1.414, 0.707, 2.121]\n"
    "  -> exp: [4.113, 2.028, 8.346]  -> normalize -> [0.284, 0.140, 0.576]\n"
    "\n"
    "row 3: [2.121, 0.707, 2.828]\n"
    "  -> exp: [8.346, 2.028, 16.91]  -> normalize -> [0.306, 0.074, 0.620]", fsz=15)

_code(s, 2.5, 5.3, 8.3, 1.3,
    "         tok1    tok2    tok3\n"
    "A  =  | 0.401   0.198   0.401 |   row sum = 1\n"
    "      | 0.284   0.140   0.576 |   \"how much of each\n"
    "      | 0.306   0.074   0.620 |    Value to blend\"", fsz=15)

# ── 13. 出力 ──
s = _content("手計算: 出力 Y = AV")

tb = _tb(s, 0.8, 1.2, 8, 0.4)
tf = tb.text_frame; p = tf.paragraphs[0]
_run(p, "Step 5: Valueの重み付き和 = 情報のブレンド", sz=19, bold=True, color=NEAR_BLACK)

add_eq(s, r"Y = AV", 0.8, 1.7, 1.8, 0.45, fontsize=22)

_code(s, 0.8, 2.3, 11.5, 2.5,
    "y1 = 0.401*[2,0] + 0.198*[0,1] + 0.401*[2,1] = [1.604, 0.599]\n"
    "\n"
    "y2 = 0.284*[2,0] + 0.140*[0,1] + 0.576*[2,1] = [1.720, 0.716]\n"
    "\n"
    "y3 = 0.306*[2,0] + 0.074*[0,1] + 0.620*[2,1] = [1.851, 0.694]", fsz=16)

_keypoint(s, 0.8, 5.2, 11.5, 0.9,
    "Attention重み(a_ij)に従って他トークンのValueをブレンドし、自分の表現を更新 = 「文脈化」")

# ── 14. 因果的Self-Attention ──
s = _content("因果的 Self-Attention (Decoder用)")

tb = _tb(s, 0.8, 1.2, 11.5, 0.6)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "「未来トークンを見てはいけない」 → スコアの未来部分を ", sz=19, color=DARK_GRAY)
_run(p, "-inf", sz=19, bold=True, color=RGBColor(0xDC,0x35,0x45), name=MONO)
_run(p, " にして softmax", sz=19, color=DARK_GRAY)

_code(s, 0.8, 2.0, 5.3, 1.7,
    "Mask matrix:\n"
    " |1  0  0|  pos 1: sees {1}\n"
    " |1  1  0|  pos 2: sees {1,2}\n"
    " |1  1  1|  pos 3: sees {1,2,3}", fsz=16)

_code(s, 6.8, 2.0, 5.5, 1.7,
    "Resulting weights:\n"
    " |1.000  0      0    |\n"
    " |0.670  0.330  0    |\n"
    " |0.306  0.074  0.620|", fsz=16)

_code(s, 0.8, 4.1, 11.5, 1.4,
    "Output:\n"
    "  y1 = [2, 0]           <- sees only self\n"
    "  y2 = [1.340, 0.330]   <- sees tok 1,2\n"
    "  y3 = [1.851, 0.694]   <- sees all (last pos)", fsz=16)

_keypoint(s, 0.8, 5.8, 11.5, 0.7,
    "未来情報の遮断が、重み行列のゼロとして現れる")

# ── 15. Cross-Attention ──
s = _content("Cross-Attention (Encoder-Decoder間)")

tb = _tb(s, 0.8, 1.2, 11.5, 0.6)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "Self-Attentionと式の形は全く同じ。", sz=19, color=DARK_GRAY)
_run(p, "Q, K, V の出どころが違うだけ。", sz=19, bold=True, color=PRIMARY)

# 3列テキスト (色付きバー付き)
cross_items = [
    ("Query",  "Decoder側 (前のdecoder層の出力)", PRIMARY),
    ("Key",    "Encoder出力 (memory)",           ACCENT),
    ("Value",  "Encoder出力 (memory)",           WARN),
]
for i, (name, desc, color) in enumerate(cross_items):
    x = 0.8 + i * 4.2
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Pt(4), Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    tb = _tb(s, x + 0.15, 2.2, 3.5, 1.2)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; _run(p, name, sz=22, bold=True, color=color)
    _para(tf, desc, sz=15, color=DARK_GRAY, sb=Pt(8))

add_eq(s,
    r"\mathrm{CrossAttn} = \mathrm{softmax}\!\left(\frac{Q_{\mathrm{dec}}\,K_{\mathrm{enc}}^\top}{\sqrt{d_k}}\right)V_{\mathrm{enc}}",
    1.5, 3.8, 10.3, 1.0, fontsize=24)

_keypoint(s, 0.8, 5.2, 11.5, 0.7,
    "Decoderトークンが「Encoderメモリ(画像+入力テキスト)」のどこを読むかを決める仕組み")

# ── 16. 学習 ──
s = _content("重み行列はどう学習されるのか")

tb = _tb(s, 0.8, 1.2, 11.5, 1.5)
tf = tb.text_frame; tf.word_wrap = True
_bullets(tf, [
    "モデルが次トークン確率を出す",
    "正解との誤差(クロスエントロピー)を計算",
    "誤差逆伝播で勾配が W_Q, W_K, W_V に流れ、更新される",
], sz=18)

add_eq(s,
    r"\frac{\partial L}{\partial s_{ij}} = \alpha_{ij} \cdot \left( g_i \cdot (v_j - y_i) \right)",
    2.0, 3.1, 9.3, 0.9, fontsize=26)

tb2 = _tb(s, 0.8, 4.3, 11.5, 1.5)
tf2 = tb2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_run(p, "この式の直感:", sz=19, bold=True, color=NEAR_BLACK)
_bullets(tf2, [
    ("v_j が損失を下げる方向に寄与 → ", "s_ij を上げるよう勾配が働く"),
    ("すると q_i と k_j の内積が増える ", "(相性が上がる)"),
], sz=18)

_keypoint(s, 0.8, 6.0, 11.5, 0.7,
    "本質: 「欲しいValueを持つKeyとQueryを、内積が大きくなる向きに整列させる」")

# ── 17. Multi-Head ──
s = _content("Multi-Head Attention")

add_eq(s,
    r"\mathrm{MultiHead}(Q,K,V) = \mathrm{Concat}(\mathrm{head}_1, \dots, \mathrm{head}_h) \, W^O",
    0.8, 1.2, 11.5, 0.65, fontsize=23)
add_eq(s,
    r"\mathrm{head}_i = \mathrm{Attention}(Q W_i^Q,\; K W_i^K,\; V W_i^V)",
    0.8, 1.9, 11.5, 0.65, fontsize=23)

tb = _tb(s, 0.8, 2.8, 11.5, 0.5)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "同じ入力から複数組の(Q,K,V)を作って並列にattentionし、結合して出力", sz=19, color=DARK_GRAY)

tb2 = _tb(s, 0.8, 3.5, 11.5, 0.4)
tf2 = tb2.text_frame; p = tf2.paragraphs[0]
_run(p, "なぜ必要か？", sz=20, bold=True, color=NEAR_BLACK)

tb3 = _tb(s, 0.8, 4.0, 11.5, 2.5)
tf3 = tb3.text_frame; tf3.word_wrap = True
p = tf3.paragraphs[0]
_run(p, "「同じトークン列でも、見るべき関係は複数ある」から", sz=19, bold=True, color=PRIMARY)
_bullets(tf3, [
    ("言語: ", "係り受け, 照応, 語彙的類似, 文脈上の役割"),
    ("画像パッチ: ", "形状の連続性, 同一物体の別パーツ, テクスチャ"),
], sz=18)
_para(tf3, "1ヘッドだけだと平均化が起きて、複数の関係を同時に捉えられない", sz=17, color=MID_GRAY, sb=Pt(12))

# ── 18. 位置埋め込み ──
s = _content("位置埋め込み (Positional Encoding)")

tb = _tb(s, 0.8, 1.2, 11.5, 1.0)
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_run(p, "Self-Attention単体は入力の順序を区別できない", sz=21, bold=True, color=PRIMARY)
_para(tf, "(並べ替えても出力が同じように並べ替わるだけ)", sz=17, color=MID_GRAY, sb=Pt(6))

add_eq(s,
    r"\mathrm{input} = \mathrm{token\;embedding} + \mathrm{positional\;encoding}",
    1.5, 2.8, 10.3, 0.8, fontsize=26)

tb2 = _tb(s, 0.8, 4.0, 11.5, 2.5)
tf2 = tb2.text_frame; tf2.word_wrap = True
_bullets(tf2, [
    "Transformerは畳み込みも再帰もないので、順序情報を明示的に注入する必要がある",
    "語順が意味を変えるタスクでは、位置を入れないと区別不可能",
    ("ViTでも同じ: ", "画像パッチに位置埋め込みを足してからTransformerに入れる"),
], sz=18)

# ── 19. RT-2 / PaLI-X ──
s = _content("RT-2 / PaLI-X での Self-Attention")

# Encoder列
bar_e = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Pt(4), Inches(2.8))
bar_e.fill.solid(); bar_e.fill.fore_color.rgb = PRIMARY; bar_e.line.fill.background()
tb_e = _tb(s, 1.0, 1.2, 5.3, 2.8)
tf_e = tb_e.text_frame; tf_e.word_wrap = True
p = tf_e.paragraphs[0]; _run(p, "Encoder", sz=22, bold=True, color=PRIMARY)
_para(tf_e, "", sz=4)
_bullets(tf_e, [
    "画像 → ViT → 画像トークン列",
    "テキスト → token emb → テキストトークン列",
    "連結して1本の系列に",
], sz=16)
_para(tf_e, "", sz=4)
p2 = _para(tf_e, "", sz=16)
_run(p2, "双方向Self-Attention ", sz=16, bold=True, color=PRIMARY)
_run(p2, "(マスクなし)", sz=15, color=MID_GRAY)

# Decoder列
bar_d = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.2), Pt(4), Inches(2.8))
bar_d.fill.solid(); bar_d.fill.fore_color.rgb = WARN; bar_d.line.fill.background()
tb_d = _tb(s, 7.2, 1.2, 5.3, 2.8)
tf_d = tb_d.text_frame; tf_d.word_wrap = True
p = tf_d.paragraphs[0]; _run(p, "Decoder", sz=22, bold=True, color=WARN)
_para(tf_d, "", sz=4)
_bullets(tf_d, [
    ("① 因果的Self-Attention ", "生成済みだけ見る"),
    ("② Cross-Attention ", "Q=dec, K/V=enc出力"),
], sz=16)

tb_rt = _tb(s, 0.8, 4.4, 11.5, 0.4)
tf_rt = tb_rt.text_frame; p = tf_rt.paragraphs[0]
_run(p, "RT-2 での追加ポイント", sz=20, bold=True, color=NEAR_BLACK)

tb_rt2 = _tb(s, 0.8, 4.9, 11.5, 2.0)
tf_rt2 = tb_rt2.text_frame; tf_rt2.word_wrap = True
_bullets(tf_rt2, [
    "Encoderは(画像+命令文)を双方向Self-Attentionで文脈化",
    "Decoderは因果的Self-Attention + Cross-Attentionで出力列を生成",
    ("核心: ", "出力列が「自然言語」ではなく「行動トークン」になっているだけ"),
], sz=18)

# ── 20. まとめ ──
s = _content("まとめ")

items = [
    ("Self-Attention",   "各トークンが列全体から必要情報を加重和で取り込み、文脈化された表現を作る"),
    ("Q, K, V",          "同じ入力から別々の線形変換で「探す軸」「タグ付けの軸」「中身の軸」を分離"),
    ("なぜ内積?",        "方向の一致度を測れる + 行列積で全ペア一括 + W_Q,W_Kの学習で意味が生まれる"),
    ("sqrt(d_k)",        "内積の分散がd_kに比例するのを正規化し、softmaxの飽和を防ぐ"),
    ("softmax → 加重和", "相性スコア → 確率分布 → Valueのブレンド"),
    ("因果マスク",       "未来トークンのスコアを -inf にして情報漏洩を防ぐ"),
    ("Cross-Attention",  "Q=decoder, K/V=encoder。Decoderが入力全体を参照する仕組み"),
    ("Multi-Head",       "複数の関係を同時に捉える (1ヘッドでは平均化で潰れる)"),
    ("位置埋め込み",     "Attention単体は順序を持てないので、位置情報を明示的に注入"),
    ("RT-2",             "VLMのAttention構造はそのまま。出力が行動トークンになっただけ"),
]

for i, (term, desc) in enumerate(items):
    y = 1.15 + i * 0.6
    # 左にアクセント丸
    _circle_num(s, 1.2, y + 0.2, i+1, color=PRIMARY if i < 5 else ACCENT)
    tb = _tb(s, 1.6, y, 11.0, 0.55)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _run(p, term, sz=15, bold=True, color=NEAR_BLACK)
    _run(p, "   " + desc, sz=14, color=DARK_GRAY)

# ── 保存 ──
output_path = "/home/raito/competition/self_attention_explained.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Generated {_eq_counter} equation images, {_slide_num} slides")
